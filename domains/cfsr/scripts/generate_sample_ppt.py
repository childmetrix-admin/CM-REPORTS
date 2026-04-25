"""
Generate a sample CFSR presentation using Kurt's template.
Creates a PPT with placeholder charts to demonstrate the layout.

Run from repo root:
  python domains/cfsr/scripts/generate_sample_ppt.py md 2026_02
"""

import sys
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont

# Configuration
BLOB_BASE = "https://stchildmetrixprod.blob.core.windows.net/processed"


def create_placeholder_image(width_px, height_px, text, color="#4472C4", output_path=None):
    """Create a placeholder image with text."""
    img = Image.new('RGB', (width_px, height_px), color='white')
    draw = ImageDraw.Draw(img)
    
    # Draw border
    border_color = color
    draw.rectangle([0, 0, width_px-1, height_px-1], outline=border_color, width=3)
    
    # Draw diagonal lines for placeholder effect
    for i in range(0, width_px + height_px, 50):
        draw.line([(i, 0), (0, i)], fill='#E5E5E5', width=1)
        draw.line([(width_px - i, height_px), (width_px, height_px - i)], fill='#E5E5E5', width=1)
    
    # Draw center rectangle for text
    rect_margin = 50
    draw.rectangle(
        [rect_margin, height_px//2 - 40, width_px - rect_margin, height_px//2 + 40],
        fill='white',
        outline=border_color,
        width=2
    )
    
    # Add text
    try:
        font = ImageFont.truetype("arial.ttf", 24)
    except:
        font = ImageFont.load_default()
    
    # Center text
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    x = (width_px - text_width) // 2
    y = (height_px - text_height) // 2
    draw.text((x, y), text, fill=border_color, font=font)
    
    if output_path:
        img.save(output_path)
    return img


def generate_sample_presentation(state: str, period: str):
    """Generate a sample CFSR presentation."""
    root = Path(__file__).resolve().parents[3]
    template_path = root / "domains" / "cfsr" / "templates" / "cfsr-presentation-template.pptx"
    
    if not template_path.exists():
        print(f"Template not found: {template_path}")
        return None
    
    state_upper = state.upper()
    state_lower = state.lower()
    state_name = {"md": "Maryland", "ky": "Kentucky"}.get(state_lower, state_upper)
    
    # Create output directory
    output_dir = root / "states" / state_lower / "cfsr" / "presentations" / period
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Create screenshots directory
    screenshots_dir = output_dir / "screenshots"
    screenshots_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"Generating sample presentation for {state_name} - {period}")
    
    # Create placeholder images
    print("Creating placeholder chart images...")
    
    # Wide image for Summary App (Top Banner: 12.89 x 5.38 inches at 150 DPI)
    summary_img_path = screenshots_dir / f"{state_lower}_summary_app_{period}.png"
    create_placeholder_image(
        1934, 807,
        f"CFSR Summary Dashboard - {state_name}",
        "#0e9ba4",
        summary_img_path
    )
    print(f"  Created: {summary_img_path.name}")
    
    # Panel images for RSP/Observed/Indicators (8.67 x 7.05 inches at 150 DPI)
    panel_images = [
        (f"{state_lower}_rsp_overview_{period}.png", f"RSP Performance Overview - {state_name}"),
        (f"{state_lower}_observed_overview_{period}.png", f"Observed Performance Overview - {state_name}"),
        (f"{state_lower}_maltreatment_in_care_{period}.png", "Maltreatment in Foster Care - By State"),
        (f"{state_lower}_maltreatment_recurrence_{period}.png", "Maltreatment Recurrence - By State"),
        (f"{state_lower}_entry_rate_{period}.png", "Foster Care Entry Rate - By State"),
        (f"{state_lower}_perm12_entries_{period}.png", "Permanency 12mo (Entries) - By State"),
        (f"{state_lower}_perm12_12_23_{period}.png", "Permanency 12mo (12-23mo) - By State"),
        (f"{state_lower}_perm12_24_{period}.png", "Permanency 12mo (24+mo) - By State"),
        (f"{state_lower}_reentry_{period}.png", "Reentry to Foster Care - By State"),
        (f"{state_lower}_placement_stability_{period}.png", "Placement Stability - By State"),
    ]
    
    for filename, text in panel_images:
        img_path = screenshots_dir / filename
        create_placeholder_image(1300, 1058, text, "#4472C4", img_path)
        print(f"  Created: {filename}")
    
    # Load template
    print("Loading template...")
    prs = Presentation(str(template_path))
    
    # Clear existing slides (template may have sample content)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # Get layouts
    layouts = {layout.name: layout for layout in prs.slide_layouts}
    print(f"Available layouts: {list(layouts.keys())}")
    
    # Slide 1: Title Slide
    print("Adding slides...")
    slide = prs.slides.add_slide(layouts["Title Slide"])
    title = slide.shapes.title
    title.text = f"{state_name} CFSR Profile"
    
    subtitle = slide.placeholders[1]
    period_display = f"{'February' if period.endswith('02') else 'August'} {period[:4]}"
    subtitle.text = f"Data Profile Period: {period_display}\nGenerated: {datetime.now().strftime('%B %d, %Y')}"
    
    # Slide 2: CFSR Round 4 Profile
    slide = prs.slides.add_slide(layouts["Title and Content"])
    slide.shapes.title.text = "CFSR Round 4 Profile"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Children's Bureau provides CFSR Round 4 Data Profiles every 6 months"
    p = tf.add_paragraph()
    p.text = "Shows your state's risk-standardized performance (RSP) and observed performance"
    p = tf.add_paragraph()
    p.text = "RSP is observed performance but with risk-adjustment"
    p = tf.add_paragraph()
    p.text = "RSP is compared to national performance to determine if statistically better, worse, or no different"
    
    # Slide 3: Section Header - Performance Summary
    slide = prs.slides.add_slide(layouts["Section Header"])
    slide.shapes.title.text = "CFSR Performance Summary"
    
    # Slide 4: Summary App (Top Banner with Picture)
    slide = prs.slides.add_slide(layouts["Top Banner with Picture"])
    slide.shapes.title.text = "Overall Performance Summary"
    pic_placeholder = slide.placeholders[1]
    pic_placeholder.insert_picture(str(summary_img_path))
    
    # Slide 5: RSP Overview (Side Panel with Picture)
    slide = prs.slides.add_slide(layouts["Side Panel with Picture"])
    slide.shapes.title.text = "Risk-Standardized Performance Overview"
    pic_placeholder = slide.placeholders[1]
    pic_placeholder.insert_picture(str(screenshots_dir / f"{state_lower}_rsp_overview_{period}.png"))
    
    # Slide 6: Observed Overview (Side Panel with Picture)
    slide = prs.slides.add_slide(layouts["Side Panel with Picture"])
    slide.shapes.title.text = "Observed Performance Overview"
    pic_placeholder = slide.placeholders[1]
    pic_placeholder.insert_picture(str(screenshots_dir / f"{state_lower}_observed_overview_{period}.png"))
    
    # Slide 7: Section Header - Individual Indicators
    slide = prs.slides.add_slide(layouts["Section Header"])
    slide.shapes.title.text = "Individual Indicators"
    
    # Indicator slides
    indicators = [
        ("Maltreatment in foster care", "maltreatment_in_care", 
         ["Performance: 10.04 per 100,000 days", "Ranks 45 of 51 reporting states", "0.97 above national standard of 9.07", "Performance below national standard"]),
        ("Maltreatment recurrence", "maltreatment_recurrence",
         ["Performance: 19.8%", "Ranks 28 of 50 reporting states", "10.1% above national standard of 9.7%", "Performance below national standard"]),
        ("Foster care entry rate", "entry_rate",
         ["Performance: 1.07 per 1,000", "Ranks 2 of 52 reporting states", "Based on 12,052 children"]),
        ("Permanency 12mo (entries)", "perm12_entries",
         ["Performance: 34.5%", "Ranks 42 of 49 reporting states", "0.7% below national standard of 35.2%", "Performance below national standard"]),
        ("Permanency 12mo (12-23mo)", "perm12_12_23",
         ["Performance: 64.3%", "Ranks 44 of 49 reporting states", "20.5% above national standard of 43.8%", "Performance meets national standard"]),
        ("Permanency 12mo (24+mo)", "perm12_24",
         ["Performance: 16.7%", "Ranks 36 of 49 reporting states", "20.6% below national standard of 37.3%", "Performance below national standard"]),
        ("Reentry to foster care", "reentry",
         ["Performance: 5.6%", "Ranks 13 of 44 reporting states", "At national standard of 5.6%"]),
        ("Placement stability", "placement_stability",
         ["Performance: 3.9 per 1,000 days", "Ranks 26 of 49 reporting states", "0.58 below national standard of 4.48", "Performance meets national standard"]),
    ]
    
    for title_text, img_stem, bullets in indicators:
        slide = prs.slides.add_slide(layouts["Side Panel with Picture"])
        slide.shapes.title.text = title_text
        
        # Add image
        pic_placeholder = slide.placeholders[1]
        img_path = screenshots_dir / f"{state_lower}_{img_stem}_{period}.png"
        pic_placeholder.insert_picture(str(img_path))
        
        # Add talking points
        text_placeholder = slide.placeholders[2]
        tf = text_placeholder.text_frame
        tf.text = bullets[0]
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
    
    # Slide: Closing/Summary
    slide = prs.slides.add_slide(layouts["Title and Content"])
    slide.shapes.title.text = f"Summary — {state_name}"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = f"CFSR Round 4 Data Profile: {period_display}"
    p = tf.add_paragraph()
    p.text = "Meets/exceeds national standard (3): Entry rate, Perm 12mo (12-23mo), Placement stability"
    p = tf.add_paragraph()
    p.text = "Below national standard (5): Maltreatment in care, Recurrence, Perm 12mo entries, Perm 12mo (24+mo), Reentry"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Review each indicator slide with agency leadership before external distribution."
    p = tf.add_paragraph()
    p.text = "Source: Children's Bureau CFSR Round 4 Data Profile"
    p = tf.add_paragraph()
    p.text = "Contact: kurt@childmetrix.com"
    
    # Save presentation
    output_path = output_dir / f"{state_upper}_CFSR_Presentation_{period}.pptx"
    prs.save(str(output_path))
    print(f"\nPresentation saved: {output_path}")
    
    return str(output_path)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python generate_sample_ppt.py <state> <period>")
        print("Example: python generate_sample_ppt.py md 2026_02")
        sys.exit(1)
    
    state = sys.argv[1]
    period = sys.argv[2]
    
    result = generate_sample_presentation(state, period)
    if result:
        print(f"\nSuccess! Open the file to see the presentation.")
    else:
        sys.exit(1)
