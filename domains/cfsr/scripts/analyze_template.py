"""Analyze Kurt's PPT template to understand structure and placeholders."""

from pptx import Presentation
from pptx.util import Inches, Pt
import sys
from pathlib import Path

def analyze_template(template_path: str) -> None:
    prs = Presentation(template_path)

    print(f"Slide dimensions: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
    print(f"Total slides: {len(prs.slides)}")
    print()

    # List all slide layouts available
    print("=== SLIDE LAYOUTS ===")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  Layout {i}: {layout.name}")
    print()

    # Analyze each slide
    print("=== SLIDE ANALYSIS ===")
    for idx, slide in enumerate(prs.slides):
        print()
        print(f"--- Slide {idx + 1} ---")
        print(f"  Layout: {slide.slide_layout.name}")
        
        for shape in slide.shapes:
            shape_info = f"  Shape: {shape.shape_type}, "
            if hasattr(shape, "name"):
                shape_info += f'name="{shape.name}", '
            if hasattr(shape, "left"):
                shape_info += f"pos=({shape.left.inches:.2f}, {shape.top.inches:.2f}), "
                shape_info += f"size=({shape.width.inches:.2f} x {shape.height.inches:.2f})"
            print(shape_info)
            
            # If it's a text shape, show text content
            if shape.has_text_frame:
                text = shape.text_frame.text[:150].replace("\n", " | ")
                if text.strip():
                    print(f'    Text: "{text}"')
            
            # If it's a picture placeholder or placeholder
            if shape.is_placeholder:
                ph = shape.placeholder_format
                print(f"    Placeholder: idx={ph.idx}, type={ph.type}")


if __name__ == "__main__":
    root = Path(__file__).resolve().parents[3]
    template_path = root / "MD_CFSR_Presentation_2026_02-Kurt.pptx"
    
    if not template_path.exists():
        print(f"Template not found: {template_path}")
        sys.exit(1)
    
    analyze_template(str(template_path))
